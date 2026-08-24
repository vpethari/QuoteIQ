"""Generate QuoteIQ POC overview PowerPoint. Run: py docs/build_quoteiq_overview_pptx.py"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
NAVY_MID = RGBColor(0x14, 0x3D, 0x66)
GREEN = RGBColor(0x7A, 0xC1, 0x43)
ORANGE = RGBColor(0xE3, 0x6C, 0x0A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x3A, 0x4A, 0x5C)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC5, 0xD0, 0xDC)
TEAL = RGBColor(0x1A, 0x7A, 0x8C)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_run(run, text: str, size: int, color: RGBColor, bold: bool = False) -> None:
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"


def add_text_box(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(_ensure_run(p), text, size, color, bold)
    return box


def _ensure_run(paragraph):
    if paragraph.runs:
        return paragraph.runs[0]
    return paragraph.add_run()


def add_box_text(shape, lines: list[tuple[str, int, RGBColor, bool]], align=PP_ALIGN.LEFT) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    first = True
    for text, size, color, bold in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(2)
        run = _ensure_run(p)
        set_run(run, text, size, color, bold)


def rect(slide, left, top, width, height, fill: RGBColor, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def pill(slide, left, top, width, height, fill: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def footer(slide, page: str) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.28), SLIDE_W, Inches(0.22))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_text_box(slide, Inches(0.4), Inches(7.28), Inches(8), Inches(0.22), "QuoteIQ  ·  Proof of Concept  ·  Confidential", 10, WHITE)
    add_text_box(slide, Inches(11.6), Inches(7.28), Inches(1.4), Inches(0.22), page, 10, WHITE, align=PP_ALIGN.RIGHT)


def header_bar(slide, title: str, subtitle: str | None = None) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.05), SLIDE_W, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN
    accent.line.fill.background()
    add_text_box(slide, Inches(0.45), Inches(0.18), Inches(11), Inches(0.45), title, 26, WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.45), Inches(0.62), Inches(12), Inches(0.32), subtitle, 13, GREEN)


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_bullet_card(slide, left, top, width, height, title, bullets, accent=GREEN):
    card = rect(slide, left, top, width, height, WHITE, LINE)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    add_text_box(slide, left + Inches(0.22), top + Inches(0.1), width - Inches(0.3), Inches(0.35), title, 15, NAVY, bold=True)
    box = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.45), width - Inches(0.35), height - Inches(0.55))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        run = _ensure_run(p)
        set_run(run, "•  " + item, 12, SLATE)
    return card


def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- 1 Title ---
    s = blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), SLIDE_H)
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN
    accent.line.fill.background()
    add_text_box(s, Inches(0.7), Inches(1.9), Inches(11), Inches(0.4), "ATKORE  ·  PROOF OF CONCEPT", 14, GREEN, bold=True)
    add_text_box(s, Inches(0.7), Inches(2.35), Inches(12), Inches(1.1), "QuoteIQ", 54, WHITE, bold=True)
    add_text_box(
        s,
        Inches(0.7),
        Inches(3.5),
        Inches(11.5),
        Inches(0.9),
        "Intelligent product matching that turns messy RFQs into CPQ-ready quotes",
        20,
        WHITE,
    )
    add_text_box(
        s,
        Inches(0.7),
        Inches(4.55),
        Inches(11),
        Inches(0.8),
        "Executive summary  ·  Logical flow  ·  Architecture  ·  Azure enterprise path",
        16,
        RGBColor(0xA8, 0xC0, 0xD4),
    )
    add_text_box(s, Inches(0.7), Inches(6.6), Inches(10), Inches(0.35), "POC overview  |  August 2026", 13, GREEN)

    # --- 2 Agenda ---
    s = blank_slide(prs)
    header_bar(s, "Agenda", "Four views of QuoteIQ for leadership and engineering")
    items = [
        ("01", "Executive summary", "What QuoteIQ does, who it helps, and the POC capabilities that are live today."),
        ("02", "Logical diagram", "Business flow from customer RFQ to catalog match, human review, and CSV export."),
        ("03", "Architecture diagram", "Current POC stack: React UI, FastAPI matching engine, Excel catalog, optional Azure OpenAI."),
        ("04", "Azure & enterprise", "Target cloud landing zone, identity, security, operations, and how the POC becomes a product."),
    ]
    for i, (num, title, body) in enumerate(items):
        top = Inches(1.4) + Inches(i * 1.35)
        n = pill(s, Inches(0.5), top + Inches(0.25), Inches(0.7), Inches(0.55), ORANGE if i == 3 else GREEN)
        add_box_text(n, [(num, 16, WHITE, True)], PP_ALIGN.CENTER)
        card = rect(s, Inches(1.45), top, Inches(11.3), Inches(1.2), WHITE, LINE)
        add_text_box(s, Inches(1.7), top + Inches(0.18), Inches(10.8), Inches(0.35), title, 18, NAVY, bold=True)
        add_text_box(s, Inches(1.7), top + Inches(0.55), Inches(10.8), Inches(0.5), body, 14, SLATE)
    footer(s, "2")

    # --- 3 Exec problem / value ---
    s = blank_slide(prs)
    header_bar(s, "1. Executive summary", "The problem QuoteIQ solves")
    add_bullet_card(
        s,
        Inches(0.4),
        Inches(1.35),
        Inches(6.1),
        Inches(5.55),
        "Business problem",
        [
            "Customer quotes arrive as Excel with descriptions, abbreviations, or mixed part numbers — not a clean Atkore SKU list.",
            "Inside sales spends time hunting the catalog, guessing among duplicate descriptions, and risking the wrong part on a quote.",
            "Family / parent catalog IDs must never be sold; official catalog numbers must be the only output.",
            "Downstream CPQ and agents need a stable CSV, not a rewritten customer workbook.",
        ],
        ORANGE,
    )
    add_bullet_card(
        s,
        Inches(6.75),
        Inches(1.35),
        Inches(6.15),
        Inches(5.55),
        "POC value",
        [
            "Upload an Excel RFQ; QuoteIQ extracts line items (Name / Description + Qty; optional Part Number).",
            "Match against the Atkore product catalog using identifiers first, then description scoring.",
            "Ambiguous or duplicate catalog descriptions are flagged for review — the engine does not invent a winner.",
            "Optional Azure OpenAI only re-ranks existing candidates; it cannot invent part numbers.",
            "Download a CPQ-ready CSV with status, score, reasons, and top candidates.",
        ],
        GREEN,
    )
    footer(s, "3")

    # --- 4 Features ---
    s = blank_slide(prs)
    header_bar(s, "1. Executive summary — POC features", "What is in the application today")
    features = [
        ("Quote ingest", "Excel upload (xlsx). Header aliases for Name/Description, Qty, optional PN. Blank Name allowed if a part number is present."),
        ("Identifier match", "Exact Salsify ID (keep NA1- prefix), then official catalog number. Name cell can be a PN or a description."),
        ("Description match", "Normalize text, expand synonyms (e.g. LTG → LIGHTING), score exact / token / fuzzy / attributes."),
        ("Catalog rules", "Sellable products only. Family/parent Salsify IDs never match. Output is always official catalog number."),
        ("Confidence & review", "EXACT / HIGH / REVIEW REQUIRED / NO MATCH. Duplicate descriptions stay in review with all candidates shown."),
        ("Optional AI layer", "Azure OpenAI judges top N candidates only. Off by default. Missing config returns a clear unavailable message."),
        ("Results UI", "Summary counts, line table, expandable candidates, AI toggle, CSV download — no source Excel rewrite."),
        ("APIs", "Health, preview, quote match, process-to-JSON, process-to-CSV. Vite UI proxies /api in local DEV."),
    ]
    for i, (title, body) in enumerate(features):
        col = i % 2
        row = i // 2
        left = Inches(0.4) + Inches(col * 6.45)
        top = Inches(1.3) + Inches(row * 1.4)
        card = rect(s, left, top, Inches(6.25), Inches(1.28), WHITE, LINE)
        add_text_box(s, left + Inches(0.2), top + Inches(0.12), Inches(5.9), Inches(0.32), title, 14, NAVY, bold=True)
        add_text_box(s, left + Inches(0.2), top + Inches(0.46), Inches(5.9), Inches(0.72), body, 12, SLATE)
    footer(s, "4")

    # --- 5 Outcomes ---
    s = blank_slide(prs)
    header_bar(s, "1. Executive summary — outcomes", "How a processed quote is classified")
    outcomes = [
        (GREEN, "Matched", "Unique identifier hit or unique high-confidence description. Official Atkore part number is filled."),
        (ORANGE, "Review required", "Viable candidates exist but the winner is ambiguous (ties, shared descriptions, thin score gap)."),
        (RGBColor(0x8A, 0x3B, 0x3B), "No match", "Nothing reaches the match threshold. CSV part number stays blank; line is still exported."),
    ]
    for i, (color, title, body) in enumerate(outcomes):
        left = Inches(0.45) + Inches(i * 4.2)
        card = rect(s, left, Inches(1.45), Inches(3.95), Inches(2.55), WHITE, LINE)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.45), Inches(3.95), Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_text_box(s, left + Inches(0.2), Inches(1.75), Inches(3.55), Inches(0.4), title, 18, NAVY, bold=True)
        add_text_box(s, left + Inches(0.2), Inches(2.25), Inches(3.55), Inches(1.5), body, 13, SLATE)
    add_bullet_card(
        s,
        Inches(0.45),
        Inches(4.2),
        Inches(12.4),
        Inches(2.7),
        "POC boundaries (intentional)",
        [
            "Column mapping uses known header aliases — not fully arbitrary agent column names (acceptable for POC).",
            "Catalog is loaded from Excel at process time (in-memory matcher). Postgres exists in docker-compose for later catalog persistence.",
            "No login, tenant isolation, PDF RFQs, or CPQ write-back in this POC. CSV is the integration contract.",
            "AI never searches the full catalog and never invents SKUs. Local demo can run with AI off.",
        ],
        NAVY,
    )
    footer(s, "5")

    # --- 6 Logical diagram ---
    s = blank_slide(prs)
    header_bar(s, "2. Logical diagram", "Business process from RFQ to CPQ-ready output")

    steps = [
        ("Customer / agent", "Excel RFQ\nName, Qty, optional PN"),
        ("Ingest & parse", "Find headers\nNormalize line items"),
        ("Match engine", "ID first, then\ndescription scoring"),
        ("Optional AI", "Re-rank top N\ncandidates only"),
        ("Decide", "Match / Review /\nNo match"),
        ("Export", "UI + CSV for\nCPQ / agents"),
    ]
    y = Inches(1.55)
    w = Inches(1.85)
    gap = Inches(0.18)
    start = Inches(0.35)
    for i, (title, body) in enumerate(steps):
        left = start + i * (w + gap)
        fill = ORANGE if i == 5 else (TEAL if i == 3 else NAVY)
        box = rect(s, left, y, w, Inches(1.7), fill)
        add_box_text(
            box,
            [
                (title, 12, WHITE, True),
                (body, 11, WHITE, False),
            ],
            PP_ALIGN.CENTER,
        )
        if i < len(steps) - 1:
            arrow = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                left + w + Inches(0.02),
                y + Inches(0.7),
                Inches(0.14),
                Inches(0.28),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GREEN
            arrow.line.fill.background()

    # Catalog + rules row
    cat = rect(s, Inches(0.35), Inches(3.5), Inches(6.2), Inches(1.55), WHITE, LINE)
    add_text_box(s, Inches(0.55), Inches(3.6), Inches(5.9), Inches(0.32), "Atkore product catalog", 14, NAVY, bold=True)
    add_text_box(
        s,
        Inches(0.55),
        Inches(3.95),
        Inches(5.9),
        Inches(0.95),
        "Excel extract of sellable products. Official part = text before “ - ” in Catalog Number – Short Description. Families (salsify:parent_id) are never quoted.",
        12,
        SLATE,
    )
    rules = rect(s, Inches(6.75), Inches(3.5), Inches(6.15), Inches(1.55), WHITE, LINE)
    add_text_box(s, Inches(6.95), Inches(3.6), Inches(5.8), Inches(0.32), "Guardrails", 14, NAVY, bold=True)
    add_text_box(
        s,
        Inches(6.95),
        Inches(3.95),
        Inches(5.8),
        Inches(0.95),
        "Do not guess on ties. Do not emit Salsify IDs or parent IDs as the matched part. Do not modify the customer workbook. Quantity is not used to pick identity.",
        12,
        SLATE,
    )

    loop = rect(s, Inches(0.35), Inches(5.25), Inches(12.55), Inches(1.7), WHITE, LINE)
    add_text_box(s, Inches(0.55), Inches(5.38), Inches(12.2), Inches(0.3), "Human in the loop", 14, NAVY, bold=True)
    add_text_box(
        s,
        Inches(0.55),
        Inches(5.75),
        Inches(12.2),
        Inches(1.0),
        "Review-required lines show scored candidates and match reasons so inside sales can pick the correct SKU. Matched and no-match lines still flow to CSV so agents get a complete, auditable file. This is the control that keeps QuoteIQ safe for quoting.",
        13,
        SLATE,
    )
    footer(s, "6")

    # --- 7 Logical matching path ---
    s = blank_slide(prs)
    header_bar(s, "2. Logical diagram — matching path", "How a single line is decided")
    path = [
        ("1", "Parse line", "Description, qty, optional PN from mapped columns."),
        ("2", "Identifier lookup", "Normalized Salsify ID, then official catalog number (Name may be an ID)."),
        ("3", "Description scoring", "If no ID hit: exact, token Dice, fuzzy, attribute overlap → 0–100."),
        ("4", "Rank & policy", "Drop weak scores. Unique high score → match. Tie / duplicate desc → review."),
        ("5", "Optional AI", "If enabled and Azure is configured: reason over top 5 catalog candidates only."),
        ("6", "Emit result", "Official PN or blank + status, %, reasons, candidate list → UI and CSV."),
    ]
    for i, (n, title, body) in enumerate(path):
        col = i % 3
        row = i // 3
        left = Inches(0.4) + Inches(col * 4.25)
        top = Inches(1.4) + Inches(row * 2.7)
        card = rect(s, left, top, Inches(4.05), Inches(2.45), WHITE, LINE)
        badge = pill(s, left + Inches(0.2), top + Inches(0.25), Inches(0.5), Inches(0.4), GREEN)
        add_box_text(badge, [(n, 14, WHITE, True)], PP_ALIGN.CENTER)
        add_text_box(s, left + Inches(0.85), top + Inches(0.28), Inches(2.95), Inches(0.4), title, 16, NAVY, bold=True)
        add_text_box(s, left + Inches(0.2), top + Inches(0.85), Inches(3.65), Inches(1.4), body, 13, SLATE)
    footer(s, "7")

    # --- 8 Architecture current ---
    s = blank_slide(prs)
    header_bar(s, "3. Architecture — current POC", "Local / docker development topology")

    # Users
    u = rect(s, Inches(0.35), Inches(1.4), Inches(2.3), Inches(1.35), ORANGE)
    add_box_text(u, [("Users", 14, WHITE, True), ("Inside sales\nLocal demo browser", 12, WHITE, False)], PP_ALIGN.CENTER)

    # Frontend
    fe = rect(s, Inches(2.9), Inches(1.4), Inches(3.15), Inches(1.35), NAVY)
    add_box_text(fe, [("Web UI", 14, WHITE, True), ("React + Vite\nlocalhost:5173\nProxy /api → :8000", 11, WHITE, False)], PP_ALIGN.CENTER)

    # API
    api = rect(s, Inches(6.3), Inches(1.35), Inches(6.6), Inches(3.55), NAVY_MID)
    add_text_box(s, Inches(6.5), Inches(1.42), Inches(6.2), Inches(0.3), "FastAPI  ·  uvicorn  ·  :8000", 13, WHITE, bold=True)

    modules = [
        (Inches(6.5), Inches(1.85), "quotes", "Excel parse\nheader aliases"),
        (Inches(8.55), Inches(1.85), "catalog", "Load products\nfrom xlsx"),
        (Inches(10.6), Inches(1.85), "matching", "Normalize, score\nID + description"),
        (Inches(6.5), Inches(3.15), "ai", "Optional Azure\nOpenAI re-rank"),
        (Inches(8.55), Inches(3.15), "output", "JSON results\nCSV schema"),
        (Inches(10.6), Inches(3.15), "config", "Thresholds\n.env settings"),
    ]
    for left, top, title, body in modules:
        b = rect(s, left, top, Inches(1.95), Inches(1.15), NAVY)
        add_box_text(b, [(title, 12, GREEN, True), (body, 10, WHITE, False)], PP_ALIGN.CENTER)

    # Data
    d1 = rect(s, Inches(0.35), Inches(3.05), Inches(2.7), Inches(1.7), WHITE, LINE)
    add_text_box(s, Inches(0.5), Inches(3.15), Inches(2.4), Inches(0.3), "Catalog Excel", 13, NAVY, bold=True)
    add_text_box(s, Inches(0.5), Inches(3.5), Inches(2.4), Inches(1.1), "data/Atkorepartsfile.xlsx\nLoaded in-memory\ninto ProductMatcher", 11, SLATE)

    d2 = rect(s, Inches(3.2), Inches(3.05), Inches(2.85), Inches(1.7), WHITE, LINE)
    add_text_box(s, Inches(3.35), Inches(3.15), Inches(2.55), Inches(0.3), "Quote Excel", 13, NAVY, bold=True)
    add_text_box(s, Inches(3.35), Inches(3.5), Inches(2.55), Inches(1.1), "Uploaded xlsx\nParsed then discarded\nNever rewritten", 11, SLATE)

    az = rect(s, Inches(0.35), Inches(5.0), Inches(5.7), Inches(1.9), WHITE, LINE)
    add_text_box(s, Inches(0.5), Inches(5.1), Inches(5.4), Inches(0.3), "Optional: Azure OpenAI", 14, NAVY, bold=True)
    add_text_box(
        s,
        Inches(0.5),
        Inches(5.45),
        Inches(5.4),
        Inches(1.3),
        "Endpoint + key + deployment in .env. Used only when AI toggle is on. Unconfigured → 503 with a user-facing “unavailable” message. Default is AI off.",
        12,
        SLATE,
    )

    pg = rect(s, Inches(6.3), Inches(5.0), Inches(6.6), Inches(1.9), WHITE, LINE)
    add_text_box(s, Inches(6.5), Inches(5.1), Inches(6.2), Inches(0.3), "Postgres 16 (docker-compose) — ready, not on the hot path", 13, NAVY, bold=True)
    add_text_box(
        s,
        Inches(6.5),
        Inches(5.45),
        Inches(6.2),
        Inches(1.3),
        "Alembic migrations exist for catalog persistence. POC matching still uses the Excel loader. This is the on-ramp to a managed Azure database without changing the matching algorithms.",
        12,
        SLATE,
    )
    footer(s, "8")

    # --- 9 Architecture APIs ---
    s = blank_slide(prs)
    header_bar(s, "3. Architecture — interfaces", "Contracts the POC already exposes")
    add_bullet_card(
        s,
        Inches(0.4),
        Inches(1.35),
        Inches(6.15),
        Inches(5.55),
        "HTTP APIs",
        [
            "GET /health — liveness and AI flag",
            "POST /api/matching/preview — one description",
            "POST /api/matching/ai-preview — same with AI",
            "POST /api/matching/quote — JSON lines or path",
            "POST /api/quote/process/results — upload → JSON summary + lines",
            "POST /api/quote/process — upload → CSV download",
            "CORS allow-list from settings; 5 MB upload cap",
        ],
        GREEN,
    )
    add_bullet_card(
        s,
        Inches(6.75),
        Inches(1.35),
        Inches(6.15),
        Inches(5.55),
        "CSV contract (do not break consumers)",
        [
            "Source file / sheet / row",
            "Requested description and quantity",
            "Matched Atkore part + description",
            "Matching percentage (no % sign)",
            "Confidence HIGH / REVIEW / LOW",
            "Match status and reason",
            "Candidate count and top candidates",
        ],
        ORANGE,
    )
    footer(s, "9")

    # --- 10 Azure target ---
    s = blank_slide(prs)
    header_bar(s, "4. Azure migration — target architecture", "Same matching core, enterprise landing zone")

    layers = [
        (Inches(0.35), "Experience", NAVY, [("SPA", "Static Web Apps\nor App Service\n+ Front Door / CDN"), ("Auth", "Entra ID (SSO)\nApp roles\nConditional Access")]),
        (Inches(3.55), "Application", TEAL, [("API", "Azure App Service\nor Container Apps\nFastAPI workers"), ("Jobs", "Catalog refresh\nqueue (Service Bus)")]),
        (Inches(6.75), "Intelligence", ORANGE, [("Match", "Existing engine\nin-process, scaled\nout by replicas"), ("AI", "Azure OpenAI\nPrivate endpoint\ncontent filters")]),
        (Inches(9.95), "Data & secrets", GREEN, [("DB", "Azure Database\nfor PostgreSQL\nFlexible Server"), ("Files", "Blob for quotes\nKey Vault secrets")]),
    ]
    for left, title, color, cards in layers:
        add_text_box(s, left, Inches(1.25), Inches(3.0), Inches(0.3), title.upper(), 11, color, bold=True)
        for i, (h, b) in enumerate(cards):
            top = Inches(1.6) + Inches(i * 1.7)
            box = rect(s, left, top, Inches(3.0), Inches(1.55), WHITE, LINE)
            stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), Inches(1.55))
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = color
            stripe.line.fill.background()
            add_text_box(s, left + Inches(0.22), top + Inches(0.15), Inches(2.65), Inches(0.3), h, 14, NAVY, bold=True)
            add_text_box(s, left + Inches(0.22), top + Inches(0.5), Inches(2.65), Inches(0.9), b, 12, SLATE)

    note = rect(s, Inches(0.35), Inches(5.15), Inches(12.6), Inches(1.8), WHITE, LINE)
    add_text_box(s, Inches(0.55), Inches(5.28), Inches(12.2), Inches(0.3), "Migration principle", 14, NAVY, bold=True)
    add_text_box(
        s,
        Inches(0.55),
        Inches(5.65),
        Inches(12.2),
        Inches(1.1),
        "Lift the existing Python matching and AI validation modules unchanged. Replace local Excel-on-disk and in-memory catalog with Blob + PostgreSQL. Put identity, networking, and observability around the same APIs so CPQ/agents keep the CSV (and JSON) contracts.",
        13,
        SLATE,
    )
    footer(s, "10")

    # --- 11 Enterprise capabilities ---
    s = blank_slide(prs)
    header_bar(s, "4. Making QuoteIQ an enterprise app", "Capabilities beyond the POC")
    ents = [
        ("Identity & access", "Entra ID SSO, role-based access (quote user, reviewer, catalog admin). No anonymous public upload in production."),
        ("Tenancy & audit", "User, timestamp, file hash, match decision, and AI prompt version stored for every run. In-memory audit store is replaced by Postgres."),
        ("Catalog operations", "Scheduled ingest of Atkore extract; versioned catalog; no matching of family rows; admin approval before publish."),
        ("Security", "Private endpoints, Key Vault, managed identity (no keys in app settings), malware scan on upload, size/type limits, WAF on Front Door."),
        ("Reliability", "Multi-instance API, health/readiness probes, autoscale, poison-message handling for catalog jobs, RPO/RTO for Postgres."),
        ("Integration", "Keep CSV for agents; add authenticated APIs for CPQ; optional Event Grid when a quote is processed; never rewrite source files."),
        ("Responsible AI", "AI remains optional and candidate-bounded. Content filters, logging without PII in prompts where possible, human review for low confidence."),
        ("Delivery", "Dev / test / prod subscriptions, IaC (Bicep or Terraform), CI/CD, environment-specific catalogs, change control for thresholds."),
    ]
    for i, (title, body) in enumerate(ents):
        col = i % 2
        row = i // 2
        left = Inches(0.4) + Inches(col * 6.45)
        top = Inches(1.28) + Inches(row * 1.4)
        card = rect(s, left, top, Inches(6.25), Inches(1.3), WHITE, LINE)
        add_text_box(s, left + Inches(0.2), top + Inches(0.12), Inches(5.9), Inches(0.3), title, 14, NAVY, bold=True)
        add_text_box(s, left + Inches(0.2), top + Inches(0.45), Inches(5.9), Inches(0.75), body, 12, SLATE)
    footer(s, "11")

    # --- 12 Phased roadmap ---
    s = blank_slide(prs)
    header_bar(s, "4. Suggested Azure rollout", "Phased path from POC to enterprise product")
    phases = [
        ("Phase 0", "Now (POC)", GREEN, ["Local UI + API", "Excel catalog", "Alias headers", "Optional Azure OpenAI", "CSV download"]),
        ("Phase 1", "Hosted pilot", TEAL, ["App Service / ACA", "Entra ID auth", "Key Vault", "Blob uploads", "App Insights"]),
        ("Phase 2", "Data platform", ORANGE, ["Azure PostgreSQL", "Catalog pipeline", "Persistent audit", "Private OpenAI", "Non-prod + prod"]),
        ("Phase 3", "Enterprise", NAVY, ["Front Door + WAF", "RBAC / reviewers", "CPQ API", "SLA & runbooks", "Broader column mapping"]),
    ]
    for i, (phase, name, color, bullets) in enumerate(phases):
        left = Inches(0.4) + Inches(i * 3.2)
        card = rect(s, left, Inches(1.4), Inches(3.0), Inches(4.55), WHITE, LINE)
        head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.4), Inches(3.0), Inches(0.95))
        head.fill.solid()
        head.fill.fore_color.rgb = color
        head.line.fill.background()
        add_text_box(s, left + Inches(0.15), Inches(1.5), Inches(2.7), Inches(0.3), phase, 12, WHITE, bold=True)
        add_text_box(s, left + Inches(0.15), Inches(1.82), Inches(2.7), Inches(0.4), name, 16, WHITE, bold=True)
        box = s.shapes.add_textbox(left + Inches(0.2), Inches(2.5), Inches(2.6), Inches(3.2))
        tf = box.text_frame
        tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_after = Pt(10)
            set_run(_ensure_run(p), "•  " + b, 13, SLATE)
    footer(s, "12")

    # --- 13 Close ---
    s = blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), SLIDE_H)
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN
    accent.line.fill.background()
    add_text_box(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.4), "SUMMARY", 14, GREEN, bold=True)
    add_text_box(s, Inches(0.7), Inches(2.45), Inches(12), Inches(1.2), "QuoteIQ is ready as a matching POC.\nAzure makes it an enterprise quoting service.", 28, WHITE, bold=True)
    add_text_box(
        s,
        Inches(0.7),
        Inches(4.15),
        Inches(11.5),
        Inches(1.5),
        "Keep the matching rules, CSV schema, and “never invent a part number” policy.\nWrap them with identity, managed data, private AI, and operations.",
        16,
        RGBColor(0xA8, 0xC0, 0xD4),
    )
    add_text_box(s, Inches(0.7), Inches(6.4), Inches(11), Inches(0.35), "Questions and next step: Phase 1 hosted pilot", 14, GREEN)

    out = Path(__file__).resolve().parent / "QuoteIQ-POC-Overview.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    path = build()
    print(path)
